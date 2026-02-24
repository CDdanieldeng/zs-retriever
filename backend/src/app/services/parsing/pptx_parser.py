"""PPTX parser using python-pptx."""

from io import BytesIO
from typing import BinaryIO

from pptx import Presentation

from app.services.parsing.base import (
    Block,
    ImageBlock,
    Loc,
    TableBlock,
    TextBlock,
)


class PptxParser:
    """Parse PPTX files. Parent = slide."""

    def parse(self, content: bytes | BinaryIO, filename: str = "") -> list[Block]:
        """Parse PPTX content into blocks."""
        if isinstance(content, bytes):
            stream = BytesIO(content)
        else:
            stream = content
        prs = Presentation(stream)
        blocks: list[Block] = []
        for slide_num, slide in enumerate(prs.slides):
            loc = Loc(slide_num=slide_num + 1)
            text_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text_parts.append(para.text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        blocks.append(
                            TableBlock(
                                content="\n".join(rows),
                                loc=loc,
                                metadata={"shape_id": shape.shape_id},
                            )
                        )
                if shape.shape_type == 13:  # Picture
                    try:
                        img = shape.image
                        img_bytes = img.blob
                        blocks.append(
                            ImageBlock(
                                content=f"[Image slide {slide_num + 1}]",
                                loc=loc,
                                metadata={"shape_id": shape.shape_id},
                                image_bytes=img_bytes,
                            )
                        )
                    except Exception:
                        pass
            if text_parts:
                text = "\n".join(p for p in text_parts if p.strip())
                if text.strip():
                    blocks.append(TextBlock(content=text.strip(), loc=loc))
        return blocks
